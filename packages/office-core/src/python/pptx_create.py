import os

from _worker import run, WorkerError
from _pptx_common import find_layout, delete_entry
from pptx import Presentation


def main(payload):
    path = payload["file"]
    slides_spec = payload["slides"]
    template = payload.get("template")

    if os.path.exists(path):
        raise WorkerError(
            "FILE_EXISTS",
            f"{path} already exists",
            "office_create makes new files; use office_edit to change existing ones, or pass a different path.",
        )

    if template:
        try:
            prs = Presentation(template)
        except Exception as e:
            raise WorkerError(
                "FILE_OPEN",
                f"Could not open {template} as .pptx: {e}",
                "Check the path; the file must be a .pptx (not legacy .ppt).",
            )
        while len(list(prs.slides)) > 0:
            delete_entry(prs, 0)
    else:
        prs = Presentation()

    skipped = []
    for i, spec in enumerate(slides_spec):
        layout = find_layout(prs, spec["layout"])
        slide = prs.slides.add_slide(layout)
        if spec.get("title") is not None:
            if slide.shapes.title is not None:
                slide.shapes.title.text = spec["title"]
            else:
                skipped.append({"slide": i, "field": "title"})
        if spec.get("bullets"):
            placeholder = next((p for p in slide.placeholders if p.placeholder_format.idx == 1), None)
            if placeholder is not None:
                placeholder.text = "\n".join(spec["bullets"])
            else:
                skipped.append({"slide": i, "field": "bullets"})
        if spec.get("notes") is not None:
            slide.notes_slide.notes_text_frame.text = spec["notes"]

    tmp = path + ".tmp-opencode-office"
    try:
        prs.save(tmp)
        os.replace(tmp, path)
    finally:
        if os.path.exists(tmp):
            os.remove(tmp)
    return {"file": path, "slides": len(slides_spec), "skipped": skipped}


run(main)
