import io
import os
import sys

from docx import Document
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from PIL import Image as PILImage


def make_docx(path):
    doc = Document()
    doc.add_heading("Quarterly Report", level=1)
    doc.add_paragraph("Q3 revenue grew 12% year over year.")
    doc.add_heading("Regional Breakdown", level=2)
    doc.add_paragraph("EMEA led growth for the third consecutive quarter.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Revenue"
    table.cell(1, 0).text = "EMEA"
    table.cell(1, 1).text = "$4.2M"
    doc.add_paragraph("Prepared by the finance team.")
    doc.save(path)


def make_pptx(path):
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "Q3 Review"
    s1.placeholders[1].text = "Finance Team"
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Highlights"
    s2.placeholders[1].text = "Revenue up 12%\nEMEA leads growth"
    s2.notes_slide.notes_text_frame.text = "Pause here for questions."
    prs.save(path)


def make_edit_docx(path):
    doc = Document()
    doc.add_heading("Edit Playground", level=1)
    p = doc.add_paragraph("Growth was ")
    strong = p.add_run("strong")
    strong.bold = True
    p.add_run(" this quarter overall.")
    doc.add_comment(runs=p.runs[:1], text="Verify this figure", author="Reviewer", initials="RV")
    doc.add_paragraph("Delete me entirely.")
    doc.add_paragraph("Style me.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "K"
    table.cell(0, 1).text = "V"
    table.cell(1, 0).text = "alpha"
    table.cell(1, 1).text = "one"
    tracked = doc.add_paragraph("Reviewed text ")
    ins = tracked._p.makeelement(
        qn("w:ins"), {qn("w:id"): "1", qn("w:author"): "Fixture", qn("w:date"): "2026-01-01T00:00:00Z"}
    )
    run_el = tracked._p.makeelement(qn("w:r"), {})
    text_el = tracked._p.makeelement(qn("w:t"), {})
    text_el.text = "with tracked insertion"
    run_el.append(text_el)
    ins.append(run_el)
    tracked._p.append(ins)
    dele = tracked._p.makeelement(
        qn("w:del"), {qn("w:id"): "2", qn("w:author"): "Fixture", qn("w:date"): "2026-01-01T00:00:00Z"}
    )
    del_run = tracked._p.makeelement(qn("w:r"), {})
    del_text = tracked._p.makeelement(qn("w:delText"), {})
    del_text.text = "obsolete text"
    del_run.append(del_text)
    dele.append(del_run)
    tracked._p.append(dele)
    link_p = doc.add_paragraph("See ")
    rid = link_p.part.relate_to(
        "https://example.com",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True,
    )
    hyperlink = link_p._p.makeelement(qn("w:hyperlink"), {qn("r:id"): rid})
    link_run = link_p._p.makeelement(qn("w:r"), {})
    link_text = link_p._p.makeelement(qn("w:t"), {})
    link_text.text = "the appendix"
    link_run.append(link_text)
    hyperlink.append(link_run)
    link_p._p.append(hyperlink)
    link_p.add_run(" for details.")
    doc.save(path)


def make_edit_pptx(path):
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "Edit Deck"
    s1.placeholders[1].text = "v1"
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Points"
    s2.placeholders[1].text = "First point\nSecond point"
    para = s2.placeholders[1].text_frame.paragraphs[0]
    run = para.runs[0]
    run.hyperlink.address = "https://example.com"
    buf = io.BytesIO()
    PILImage.new("RGB", (64, 64), (200, 30, 30)).save(buf, format="PNG")
    buf.seek(0)
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    s3.shapes.add_picture(buf, Inches(1), Inches(1))
    chart_data = CategoryChartData()
    chart_data.categories = ["A", "B"]
    chart_data.add_series("S1", (1.0, 2.0))
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    s4.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(4), Inches(3), chart_data)
    buf2 = io.BytesIO()
    PILImage.new("RGB", (32, 32), (20, 160, 60)).save(buf2, format="PNG")
    buf2.seek(0)
    s4.shapes.add_picture(buf2, Inches(5), Inches(1))
    prs.save(path)


def make_png(path, color, size):
    PILImage.new("RGB", size, color).save(path)


if __name__ == "__main__":
    out = sys.argv[1]
    os.makedirs(out, exist_ok=True)
    make_docx(os.path.join(out, "report.docx"))
    make_pptx(os.path.join(out, "deck.pptx"))
    make_edit_docx(os.path.join(out, "edit-report.docx"))
    make_edit_pptx(os.path.join(out, "edit-deck.pptx"))
    make_png(os.path.join(out, "swap.png"), (30, 30, 200), (32, 32))
    print("ok")
