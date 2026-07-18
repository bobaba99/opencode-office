"""Test-support probe: report picture parts per slide (part name, sha256, content type)."""
import hashlib

from _worker import run
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


def main(payload):
    prs = Presentation(payload["file"])
    pictures = []
    for si, slide in enumerate(prs.slides):
        for shi, shape in enumerate(slide.shapes):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            blip = next(iter(shape._element.iter(qn("a:blip"))))
            part = slide.part.related_part(blip.get(qn("r:embed")))
            pictures.append(
                {
                    "id": f"s:{si}/sh:{shi}",
                    "part": str(part.partname),
                    "sha256": hashlib.sha256(part.blob).hexdigest(),
                    "content_type": part.content_type,
                }
            )
    return {"pictures": pictures}


run(main)
