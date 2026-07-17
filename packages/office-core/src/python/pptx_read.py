from _worker import run, WorkerError
from pptx import Presentation


def main(payload):
    path = payload["file"]
    mode = payload["mode"]
    target = payload.get("target")
    try:
        prs = Presentation(path)
    except Exception as e:
        raise WorkerError("FILE_OPEN", f"Could not open {path} as .pptx: {e}", "Check the path; the file must be a .pptx (not legacy .ppt).")
    slides = []
    for si, slide in enumerate(prs.slides):
        sid = f"s:{si}"
        if target and target != sid and not target.startswith(sid + "/"):
            continue
        title = slide.shapes.title.text if slide.shapes.title is not None else ""
        entry = {"id": sid, "title": title, "layout": slide.slide_layout.name}
        if mode == "content":
            shapes = []
            for shi, shape in enumerate(slide.shapes):
                if not shape.has_text_frame:
                    continue
                text = "\n".join(p.text for p in shape.text_frame.paragraphs)
                shapes.append({"id": f"s:{si}/sh:{shi}", "name": shape.name, "text": text})
            entry["shapes"] = shapes
            if slide.has_notes_slide:
                entry["notes"] = slide.notes_slide.notes_text_frame.text
        slides.append(entry)
    if target and not slides:
        raise WorkerError("TARGET_NOT_FOUND", f"No slide {target} in {path}", "Slide IDs come from office_read output and shift after edits; re-read to refresh them.")
    return {"format": "pptx", "mode": mode, "slides": slides}


run(main)
