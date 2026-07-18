import os
from copy import deepcopy

from _worker import run, WorkerError
from _pptx_common import find_layout, move_entry, delete_entry
from _textops import para_text, replace_in_paragraph
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


def slide_at(prs, sid):
    n = int(sid.split(":")[1].split("/")[0])
    slides = list(prs.slides)
    if n >= len(slides):
        raise WorkerError(
            "TARGET_NOT_FOUND",
            f"No slide s:{n}",
            "Slide IDs shift as the batch applies; re-read the file, or order structural ops last.",
        )
    return n, slides[n]


def shape_at(slide, sid):
    m = int(sid.split("/sh:")[1])
    shapes = list(slide.shapes)
    if m >= len(shapes):
        raise WorkerError(
            "TARGET_NOT_FOUND",
            f"No shape {sid}",
            "Shape IDs come from office_read content mode; re-read the file to refresh them.",
        )
    return shapes[m]


def replace_in_frame(shape, anchor, text, target):
    if not shape.has_text_frame:
        raise WorkerError(
            "BAD_TARGET_KIND",
            f"{target} has no text frame",
            "set_shape_text targets text shapes; office_read content mode lists them.",
        )
    frame = shape.text_frame
    total = sum(para_text(p).count(anchor) for p in frame.paragraphs)
    if total == 0:
        current = "\n".join(para_text(p) for p in frame.paragraphs)
        raise WorkerError(
            "ANCHOR_MISMATCH",
            f"Anchor not found in {target}",
            f"Shape currently reads: {current[:300]!r} — anchors cannot span paragraph breaks. The batch was rolled back — nothing was written; re-read the file and retry with corrected ops.",
        )
    if total > 1:
        raise WorkerError(
            "AMBIGUOUS_ANCHOR",
            f"Anchor occurs {total} times in {target}",
            "Extend the anchor with surrounding text until it is unique.",
        )
    for p in frame.paragraphs:
        if replace_in_paragraph(p, anchor, text):
            return


# OOXML relationships namespace. Checked by prefix rather than an enumerated attribute list
# (r:id/r:embed/r:link) because SmartArt's <dgm:relIds> uses r:dm/r:lo/r:qs/r:cs — same
# namespace, different local names — and would otherwise slip past the guard below.
R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


def copy_slide(prs, source):
    new_slide = prs.slides.add_slide(source.slide_layout)
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)
    rid_map = {}
    for rid, rel in source.part.rels.items():
        if rel.is_external:
            continue
        if "image" in rel.reltype or "media" in rel.reltype:
            rid_map[rid] = new_slide.part.relate_to(rel.target_part, rel.reltype)
    for shape in source.shapes:
        el = deepcopy(shape._element)
        for blip in el.iter(qn("a:blip")):
            embed = blip.get(qn("r:embed"))
            if embed in rid_map:
                blip.set(qn("r:embed"), rid_map[embed])
        new_slide.shapes._spTree.append(el)
    if source.has_notes_slide:
        new_slide.notes_slide.notes_text_frame.text = source.notes_slide.notes_text_frame.text

    # Fail loud on content duplicate_slide cannot copy safely: any relationship-bearing
    # attribute (chart parts, SmartArt data, OLE objects, external links, hyperlinks, ...)
    # that didn't get remapped above still points at a relationship id that doesn't exist on
    # the new slide's part. Raised before returning — apply_one/main never reach prs.save(),
    # so the batch aborts and the file on disk is untouched.
    for el in new_slide.shapes._spTree.iter():
        for name, value in el.attrib.items():
            if name.startswith(R_NS) and value not in new_slide.part.rels:
                tag = el.tag.rsplit("}", 1)[-1]
                raise WorkerError(
                    "UNSUPPORTED_SLIDE_CONTENT",
                    f"Slide contains content duplicate_slide cannot copy safely ({tag})",
                    "Charts, SmartArt, and embedded objects are not yet supported by duplicate_slide. Edit the original slide, or delete/replace the unsupported element first.",
                )
    return new_slide


def apply_one(prs, op):
    kind = op["op"]
    shape_ops = {"set_shape_text", "replace_image"}
    ident = op.get("target") or op.get("after") or ""
    if kind in shape_ops and "/sh:" not in ident:
        raise WorkerError("BAD_TARGET_KIND", f"{kind} needs a shape target (s:<n>/sh:<m>), got {ident}", "office_read content mode lists each slide's shape IDs.")
    if kind not in shape_ops and "/sh:" in ident:
        raise WorkerError("BAD_TARGET_KIND", f"{kind} needs a slide target (s:<n>), got {ident}", "Use the slide ID without the /sh:<m> suffix.")
    if kind == "set_shape_text":
        _, slide = slide_at(prs, op["target"])
        shape = shape_at(slide, op["target"])
        replace_in_frame(shape, op["anchor"], op["text"], op["target"])
        return {"op": kind, "target": op["target"]}
    if kind == "set_notes":
        _, slide = slide_at(prs, op["target"])
        slide.notes_slide.notes_text_frame.text = op["text"]
        return {"op": kind, "target": op["target"]}
    if kind == "insert_slide":
        n, _ = slide_at(prs, op["after"])
        layout = find_layout(prs, op["layout"])
        new_slide = prs.slides.add_slide(layout)
        if op.get("title") is not None and new_slide.shapes.title is not None:
            new_slide.shapes.title.text = op["title"]
        if op.get("bullets"):
            for placeholder in new_slide.placeholders:
                if placeholder.placeholder_format.idx == 1:
                    placeholder.text = "\n".join(op["bullets"])
                    break
        move_entry(prs, len(list(prs.slides)) - 1, n + 1)
        return {"op": kind, "after": op["after"]}
    if kind == "duplicate_slide":
        n, slide = slide_at(prs, op["target"])
        copy_slide(prs, slide)
        move_entry(prs, len(list(prs.slides)) - 1, n + 1)
        return {"op": kind, "target": op["target"]}
    if kind == "delete_slide":
        n, _ = slide_at(prs, op["target"])
        delete_entry(prs, n)
        return {"op": kind, "target": op["target"]}
    if kind == "move_slide":
        n, _ = slide_at(prs, op["target"])
        move_entry(prs, n, op["index"])
        return {"op": kind, "target": op["target"], "index": op["index"]}
    if kind == "replace_image":
        _, slide = slide_at(prs, op["target"])
        shape = shape_at(slide, op["target"])
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            raise WorkerError(
                "SHAPE_NOT_PICTURE",
                f"{op['target']} is not a picture shape",
                "replace_image targets picture shapes; text shapes keep their id but have no image to swap.",
            )
        try:
            image_part, rid = slide.part.get_or_add_image_part(op["image"])
        except OSError as e:
            raise WorkerError("FILE_OPEN", f"Cannot read image {op['image']}: {e}", "Check the image path exists and is readable.")
        for blip in shape._element.iter(qn("a:blip")):
            blip.set(qn("r:embed"), rid)
            break
        return {"op": kind, "target": op["target"]}
    raise WorkerError(
        "UNKNOWN_OP",
        f"Unknown pptx op: {kind}",
        "Valid ops: set_shape_text, set_notes, insert_slide, duplicate_slide, delete_slide, move_slide, replace_image.",
    )


def main(payload):
    path = payload["file"]
    try:
        prs = Presentation(path)
    except Exception as e:
        raise WorkerError("FILE_OPEN", f"Could not open {path} as .pptx: {e}", "Check the path; the file must be a .pptx (not legacy .ppt).")
    results = [apply_one(prs, op) for op in payload["operations"]]
    tmp = path + ".tmp-opencode-office"
    try:
        prs.save(tmp)
        os.replace(tmp, path)
    finally:
        if os.path.exists(tmp):
            os.remove(tmp)
    return {"applied": len(results), "results": results}


run(main)
